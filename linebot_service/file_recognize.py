import fitz  # PyMuPDF
from docx import Document
import csv
from bs4 import BeautifulSoup
from pptx import Presentation

class FileReader:
    
    @staticmethod
    def read_pdf(file_path):
        try:
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        except Exception as e:
            return f"Error reading PDF file: {e}"

    @staticmethod
    def read_word(file_path):
        try:
            doc = Document(file_path)
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
            return text
        except Exception as e:
            return f"Error reading Word file: {e}"

    @staticmethod
    def read_txt(file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return f"Error reading text file: {e}"

    @staticmethod
    def read_html(file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
                return soup.get_text()
        except Exception as e:
            return f"Error reading HTML file: {e}"

    @staticmethod
    def read_csv(file_path):
        try:
            with open(file_path, newline='', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                return "\n".join(','.join(row) for row in reader)
        except Exception as e:
            return f"Error reading CSV file: {e}"

    @staticmethod
    def read_ppt(file_path):
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"Error reading PPT file: {e}"

    def read_file(self, file_path):
        file_extension = file_path.split('.')[-1].lower()
        readers = {
            'pdf': self.read_pdf,
            'docx': self.read_word,
            'txt': self.read_txt,
            'html': self.read_html,
            'csv': self.read_csv,
            'pptx': self.read_ppt
        }
        reader = readers.get(file_extension)
        if reader:
            return reader(file_path)
        else:
            return "Unsupported file type"